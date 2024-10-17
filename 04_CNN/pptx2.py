from pptx import Presentation
import os


def extract_images_from_pptx(pptx_file_path, save_folder):
    # Load the presentation
    prs = Presentation(pptx_file_path)

    # Create the save folder if it doesn't exist
    if not os.path.exists(save_folder):
        os.makedirs(save_folder)

    # Initialize a counter for image file names
    image_count = 0

    # Iterate through all the slides
    for slide_number, slide in enumerate(prs.slides):
        # Iterate through all shapes in the slide
        for shape in slide.shapes:
            if shape.shape_type == 13:  # This is the type number for a picture
                # Access the image in the shape
                image = shape.image

                # Create an image file name
                image_filename = f"image_{slide_number+1}_{image_count+1}.png"
                image_path = os.path.join(save_folder, image_filename)

                # Save the image
                with open(image_path, 'wb') as img_file:
                    img_file.write(image.blob)

                # Increment the image counter
                image_count += 1

    print(f"Extracted {image_count} images from the presentation.")


# Example usage
pptx_file_path = r'C:\\Users\\braun\\Insync\\braunphil@gmail.com\\Google Drive\\erlangen\\lecture\\SS24\\DataScience for EM\\06_convnets.pptx'
save_folder = r'C:\\Users\\braun\\OneDrive\\Documents\\GitHub\\SS24_DataScienceForEM2\\04_CNN\\img'
extract_images_from_pptx(pptx_file_path, save_folder)
